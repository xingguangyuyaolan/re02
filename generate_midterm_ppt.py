#!/usr/bin/env python3
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
OUT_DIR = ROOT / "artifacts" / "presentations"
OUT_DIR.mkdir(parents=True, exist_ok=True)
OUT_PATH = OUT_DIR / "中期检查_基于注意力机制的多智能体强化学习无人机路径覆盖方法研究_20260419.pptx"

# Evidence assets
RUN_NEW = ROOT / "artifacts" / "qmix" / "p0_stage1_open_seed7_20260419_154937"
EVAL_DIR = RUN_NEW / "evaluation" / "20260419_192059"

IMG_LEARNING = RUN_NEW / "plots" / "learning_curves.png"
IMG_TASK = RUN_NEW / "plots" / "task_metrics.png"
IMG_BENCH = RUN_NEW / "plots" / "benchmark_metrics.png"
IMG_EVAL_PERF = EVAL_DIR / "eval_performance.png"
IMG_EVAL_TASK = EVAL_DIR / "eval_task_metrics.png"

prs = Presentation()


def add_title_slide(title: str, subtitle: str):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullets_slide(title: str, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    tf = slide.placeholders[1].text_frame
    tf.clear()
    for i, item in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.font.size = Pt(22)


def add_two_image_slide(title: str, left_img: Path, right_img: Path, caption_left: str, caption_right: str):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    top = Inches(1.2)
    h = Inches(4.8)
    if left_img.exists():
        slide.shapes.add_picture(str(left_img), Inches(0.4), top, width=Inches(6.1), height=h)
    if right_img.exists():
        slide.shapes.add_picture(str(right_img), Inches(6.8), top, width=Inches(6.1), height=h)

    tx1 = slide.shapes.add_textbox(Inches(0.4), Inches(6.1), Inches(6.1), Inches(0.6))
    tx1.text_frame.text = caption_left
    tx1.text_frame.paragraphs[0].font.size = Pt(14)

    tx2 = slide.shapes.add_textbox(Inches(6.8), Inches(6.1), Inches(6.1), Inches(0.6))
    tx2.text_frame.text = caption_right
    tx2.text_frame.paragraphs[0].font.size = Pt(14)


def add_metrics_table_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = "实验结果对比（Stage1）"

    rows, cols = 5, 6
    table = slide.shapes.add_table(rows, cols, Inches(0.35), Inches(1.3), Inches(12.6), Inches(3.2)).table

    headers = ["运行", "覆盖率", "全覆盖成功率", "出界率", "碰撞率", "总步数"]
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h

    data = [
        ["20260418_202421", "0.356", "9.62%", "93.75%", "88.94%", "12952"],
        ["20260419_113356", "0.489", "36.76%", "82.30%", "98.51%", "41439"],
        ["20260419_154937", "0.453", "44.29%", "80.09%", "93.84%", "67075"],
        ["相对早期变化", "+27.27%", "+360.57%", "-14.57%", "+5.51%", "+417.87%"],
    ]

    for r in range(1, rows):
        for c in range(cols):
            table.cell(r, c).text = data[r - 1][c]

    note = slide.shapes.add_textbox(Inches(0.45), Inches(4.8), Inches(12.2), Inches(1.5))
    ntf = note.text_frame
    ntf.text = "结论：覆盖能力与全覆盖成功率显著提升，出界率下降；碰撞率仍偏高，是下一阶段重点优化方向。"
    ntf.paragraphs[0].font.size = Pt(18)


add_title_slide(
    "中期检查汇报",
    "基于注意力机制的多智能体强化学习无人机路径覆盖方法研究\n汇报人：XXX  时间：2026-04-19",
)

add_bullets_slide(
    "研究目标与意义",
    [
        "目标：提升多无人机在复杂环境中的协同覆盖效率与稳定性",
        "问题：高碰撞率、高出界率、覆盖效率低、训练不稳定",
        "方法：Attention + QMIX + 覆盖奖励重构 + 课程学习",
        "应用：巡检、灾后搜索、农业监测等多无人机任务",
    ],
)

add_bullets_slide(
    "技术路线与系统架构",
    [
        "环境：ROS2 Jazzy + Gazebo Harmonic，多无人机并行仿真",
        "算法：Attention-QMIX（自注意力/跨智能体注意力/混合网络注意力）",
        "训练闭环：观测 -> 动作 -> 环境反馈 -> 联合价值更新",
        "课程学习：Stage1(open) -> Stage2(easy) -> Stage3(maze1)",
    ],
)

add_bullets_slide(
    "已完成工作",
    [
        "完成Attention-QMIX训练与评测主链路，实现可复现实验流程",
        "完成覆盖任务观测增强：local_map + other_agents等协同特征",
        "完成奖励重构与安全约束：边界软惩罚、碰撞冷却、低空惩罚",
        "修复关键工程问题：resume目录错误、epsilon跨阶段继承、评测参数不同步",
    ],
)

add_metrics_table_slide()

add_two_image_slide(
    "最新训练曲线证据",
    IMG_LEARNING,
    IMG_TASK,
    "左：学习曲线（learning_curves）",
    "右：任务指标（task_metrics）",
)

add_two_image_slide(
    "训练效率与评测结果",
    IMG_BENCH,
    IMG_EVAL_PERF,
    "左：训练性能指标（benchmark_metrics）",
    "右：评测性能（eval_performance）",
)

add_two_image_slide(
    "评测任务指标",
    IMG_EVAL_TASK,
    IMG_EVAL_TASK,
    "评测任务指标（collision/oob/repeat/overlap/completion）",
    "评测任务指标（同图放大展示）",
)

add_bullets_slide(
    "当前问题与原因分析",
    [
        "虽然训练阶段成功率提升，但评测全覆盖成功率仍为0，泛化稳定性不足",
        "碰撞率长期偏高，表明动作安全约束与协同避障策略仍需增强",
        "后期存在性能回落，说明奖励设计与探索衰减节奏仍需匹配优化",
        "本阶段核心成果是“可学习且可提升”，下一阶段重点是“稳定泛化”",
    ],
)

add_bullets_slide(
    "下一阶段计划",
    [
        "完成Stage2/Stage3课程训练并稳定收敛",
        "开展消融实验：无注意力/自注意力/跨智能体注意力/全注意力",
        "开展对比实验：与基础QMIX、规则法进行系统对比",
        "以出界率、碰撞率、全覆盖成功率为核心指标持续优化",
    ],
)

add_bullets_slide(
    "中期总结",
    [
        "已完成：方法搭建、训练闭环、评测闭环、关键工程修复",
        "阶段成果：全覆盖成功率与覆盖能力相较早期基线显著提升",
        "待突破点：碰撞率偏高、评测泛化稳定性不足",
        "总体判断：课题技术路线可行，后续工作目标明确",
    ],
)

prs.save(str(OUT_PATH))
print(str(OUT_PATH))
